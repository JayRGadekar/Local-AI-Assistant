import os
from typing import Optional
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from langchain_core.tools import tool

@tool
def read_local_file(filepath: str) -> str:
    """
    Read the content of a specific file from the local filesystem.
    
    Args:
        filepath: The full absolute path of the file to read.
    """
    print(f"DEBUG TOOL: read_local_file called with filepath='{filepath}'")
    return UniversalDocReader.read_file(filepath)

class UniversalDocReader:
    """
    A tool to read content from various file formats:
    - PDF (.pdf)
    - Word (.docx)
    - PowerPoint (.pptx)
    - Excel (.xlsx, .xls)
    """

    @staticmethod
    def read_file(filepath: str) -> str:
        if not os.path.exists(filepath):
            return f"Error: File not found at {filepath}"

        ext = os.path.splitext(filepath)[1].lower()

        try:
            if ext == '.pdf':
                return UniversalDocReader._read_pdf(filepath)
            elif ext == '.docx':
                return UniversalDocReader._read_docx(filepath)
            elif ext == '.pptx':
                return UniversalDocReader._read_pptx(filepath)
            elif ext in ['.xlsx', '.xls']:
                return UniversalDocReader._read_excel(filepath)
            elif ext in ['.txt', '.md', '.py', '.json']:
                return UniversalDocReader._read_text(filepath)
            else:
                return f"Error: Unsupported file extension '{ext}'"
        except Exception as e:
            return f"Error reading file: {str(e)}"

    @staticmethod
    def _read_pdf(filepath: str) -> str:
        text = ""
        with fitz.open(filepath) as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text.strip()

    @staticmethod
    def _read_docx(filepath: str) -> str:
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _read_pptx(filepath: str) -> str:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _read_excel(filepath: str) -> str:
        # Read all sheets, convert to string representation
        dfs = pd.read_excel(filepath, sheet_name=None)
        text = []
        for sheet_name, df in dfs.items():
            text.append(f"--- Sheet: {sheet_name} ---")
            text.append(df.to_string())
        return "\n".join(text)

    @staticmethod
    def _read_text(filepath: str) -> str:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
